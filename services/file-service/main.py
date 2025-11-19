from fastapi import FastAPI, UploadFile, File, Depends, HTTPException, status, BackgroundTasks
from fastapi.responses import JSONResponse
from fastapi import status as fastapi_status, Request

# Initialize FastAPI app
app = FastAPI(
    title="File Management Service",
    description="Handles file uploads, text extraction, and topic mapping",
    version="1.0.0"
)

@app.exception_handler(Exception)
async def generic_exception_handler(request: Request, exc: Exception):
    import traceback
    print("Exception:", exc)
    print(traceback.format_exc())
    return JSONResponse(
        status_code=fastapi_status.HTTP_400_BAD_REQUEST,
        content={
            "error": type(exc).__name__,
            "detail": str(exc),
            "trace": traceback.format_exc()
        },
    )
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel
from typing import List, Optional
import os
import shutil
import uuid
from datetime import datetime
from supabase import create_client, Client
import PyPDF2
from pptx import Presentation
import io

# Initialize FastAPI app
app = FastAPI(
    title="File Management Service",
    description="Handles file uploads, text extraction, and topic mapping",
    version="1.0.0"
)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Supabase client
supabase: Client = create_client(
    os.getenv("SUPABASE_URL", ""),
    os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")
)

# Security
security = HTTPBearer()

# Models
class FileResponse(BaseModel):
    id: str
    filename: str
    status: str
    uploaded_at: datetime
    error_message: Optional[str] = ""
    size: int

class TopicResponse(BaseModel):
    topic_id: str
    name: str
    doc_id: str
    page_range: Optional[str] = None

class ExtractRequest(BaseModel):
    file_id: str

# Dependency to get current user
async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    try:
        token = credentials.credentials
        user = supabase.auth.get_user(token)
        if not user:
            raise HTTPException(status_code=401, detail="Invalid authentication credentials")
        return user.user
    except Exception as e:
        raise HTTPException(status_code=401, detail=str(e))

# Helper functions
def extract_text_from_file(file_content: bytes, file_extension: str) -> str:
    text = ""
    if file_extension == '.pdf':
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
    elif file_extension in ['.ppt', '.pptx']:
        prs = Presentation(io.BytesIO(file_content))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

async def process_file_background(file_id: str, user_id: str, content: bytes, filename: str):
    try:
        # Update status to processing
        supabase.table("files").update({"status": "processing"}).eq("id", file_id).execute()
        
        # Upload to Supabase Storage
        storage_path = f"{user_id}/{file_id}/{filename}"
        supabase.storage.from_("documents").upload(
            path=storage_path,
            file=content,
            file_options={"content-type": "application/pdf" if filename.endswith('.pdf') else "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
        )
        
        # Extract text
        file_extension = os.path.splitext(filename)[1].lower()
        text = extract_text_from_file(content, file_extension)
        
        # Store extracted text in documents table
        supabase.table("documents").insert({
            "file_id": file_id,
            "text": text,
            "page_number": 1  # Simplified for now, could be per page
        }).execute()
        
        # Update status to completed
        supabase.table("files").update({
            "status": "completed",
            "storage_path": storage_path
        }).eq("id", file_id).execute()
        
    except Exception as e:
        print(f"Error processing file {file_id}: {str(e)}")
        supabase.table("files").update({
            "status": "failed",
            "error_message": str(e)
        }).eq("id", file_id).execute()

# Routes
@app.post("/files/upload", response_model=FileResponse, status_code=status.HTTP_201_CREATED)
async def upload_file(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    current_user = Depends(get_current_user)
):
    """
    Upload a PDF or PPT file for processing.
    """
    allowed_extensions = {'.pdf', '.ppt', '.pptx'}
    file_ext = os.path.splitext(file.filename)[1].lower()
    
    if file_ext not in allowed_extensions:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="File type not allowed. Please upload PDF or PPT files."
        )
    
    try:
        content = await file.read()
        file_id = str(uuid.uuid4())
        
        # Create file record
        file_data = {
            "id": file_id,
            "owner_id": current_user.id,
            "filename": file.filename,
            "status": "pending",
            "size": len(content),
            "uploaded_at": datetime.now().isoformat()
        }
        
        supabase.table("files").insert(file_data).execute()
        
        # Process in background
        background_tasks.add_task(
            process_file_background, 
            file_id, 
            current_user.id, 
            content, 
            file.filename
        )
        
        return FileResponse(**file_data)
        
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e)
        )

@app.get("/files/{file_id}/status", response_model=FileResponse)
async def get_file_status(file_id: str, current_user = Depends(get_current_user)):
    """
    Check the processing status of a file.
    """
    try:
        result = supabase.table("files").select("*").eq("id", file_id).eq("owner_id", current_user.id).execute()
        if not result.data:
            raise HTTPException(status_code=404, detail="File not found")
        return FileResponse(**result.data[0])
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/files", response_model=List[FileResponse])
async def list_files(current_user = Depends(get_current_user)):
    """
    List all uploaded files for the current user.
    """
    try:
        result = supabase.table("files").select("*").eq("owner_id", current_user.id).order("uploaded_at", desc=True).execute()
        print("Supabase result:", result)
        files = []
        for item in result.data:
            print(f"Processing file record: {item}")
            # Ensure error_message is always a string for API response
            item["error_message"] = str(item.get("error_message") or "")
            try:
                files.append(FileResponse(**item))
            except Exception as ve:
                print(f"Validation error for item: {item}\nError: {ve}")
        print(f"Returning files: {files}")
        return files
    except Exception as e:
        import traceback
        print(f"Exception in list_files: {e}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/files/{file_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_file(file_id: str, current_user = Depends(get_current_user)):
    """
    Delete a file and its associated data.
    """
    try:
        # Verify ownership
        result = supabase.table("files").select("storage_path").eq("id", file_id).eq("owner_id", current_user.id).execute()
        
        if not result.data:
            raise HTTPException(status_code=404, detail="File not found")
            
        storage_path = result.data[0].get("storage_path")
        
        # Delete from storage if path exists
        if storage_path:
            supabase.storage.from_("documents").remove([storage_path])
            
        # Delete from database (cascade will handle related records)
        supabase.table("files").delete().eq("id", file_id).execute()
        
        return None
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/health")
async def health_check():
    return {"status": "healthy", "service": "file-service"}
